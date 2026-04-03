"""
Excel → PowerPoint Automation  (Project 3)
===========================================
Reads financial_report.xlsx and produces a dark-themed CFO deck:

  Slide 1 — Cover
  Slide 2 — Revenue by Product Line    (table + clustered bar chart)
  Slide 3 — Capital Expenditure        (table + stacked bar chart)
  Slide 4 — Debt & Tax Overview        (table + dual line chart)

Usage:
    python src/excel_to_ppt.py
    python src/excel_to_ppt.py --input data/financial_report.xlsx
                               --output output/quarterly_report.pptx
                               --company "Acme Corp"
                               --period "Q4 2024 Financial Review"
"""

import argparse
from copy import deepcopy
from pathlib import Path

import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Design system
# ---------------------------------------------------------------------------

# Colours
BG          = RGBColor(0x0D, 0x1B, 0x2A)   # very dark navy — slide background
BG_HEX      = "0D1B2A"
PANEL       = RGBColor(0x12, 0x27, 0x3D)   # slightly lighter navy — chart/table bg
PANEL_HEX   = "12273D"
HDR         = RGBColor(0x1E, 0x3A, 0x5F)   # header row fill
HDR_HEX     = "1E3A5F"
ALT         = RGBColor(0x16, 0x2A, 0x46)   # alternate row fill
ALT_HEX     = "162A46"
RULE        = RGBColor(0x2C, 0x4A, 0x6E)   # subtle divider / border
RULE_HEX    = "2C4A6E"
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED       = RGBColor(0x8B, 0xA9, 0xC8)   # secondary text
ACCENT      = RGBColor(0x4F, 0xC3, 0xF7)   # light blue accent

# Series palette (chart series colours)
PALETTE = [
    RGBColor(0x4F, 0xC3, 0xF7),   # sky blue      — Product A / Maintenance
    RGBColor(0x81, 0xC7, 0x84),   # green          — Product B / Growth
    RGBColor(0xFF, 0xB7, 0x4D),   # amber          — Product C
    RGBColor(0xCE, 0x93, 0xD8),   # lavender
]
PALETTE_HEX = ["4FC3F7", "81C784", "FFB74D", "CE93D8"]

# Fonts
FONT_TITLE   = "Calibri"
FONT_BODY    = "Calibri"

# Slide geometry  (13.33" × 7.5" widescreen)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# Fixed layout zones (inches)
TITLE_L, TITLE_T, TITLE_W, TITLE_H = 0.35, 0.22, 12.63, 0.75
BAND_L,  BAND_T,  BAND_W,  BAND_H  = 0.35, 0.95, 12.63, 0.04
TABLE_L, TABLE_T, TABLE_W, TABLE_H = 0.35, 1.10,  5.90, 5.90
CHART_L, CHART_T, CHART_W, CHART_H = 6.55, 1.05,  6.43, 6.00
FOOT_L,  FOOT_T,  FOOT_W,  FOOT_H  = 0.35, 7.10, 12.63, 0.28

# ---------------------------------------------------------------------------
# XML helpers  (dark chart styling)
# ---------------------------------------------------------------------------

_A  = "http://schemas.openxmlformats.org/drawingml/2006/main"
_C  = "http://schemas.openxmlformats.org/drawingml/2006/chart"

def _solid_fill_elem(hex_color: str) -> etree._Element:
    """Return <a:solidFill><a:srgbClr val='...'></a:solidFill>."""
    fill = etree.Element(f"{{{_A}}}solidFill")
    clr  = etree.SubElement(fill, f"{{{_A}}}srgbClr")
    clr.set("val", hex_color)
    return fill

def _spPr_with_fill(hex_color: str) -> etree._Element:
    """Return <c:spPr> containing a solid fill."""
    spPr = etree.Element(f"{{{_C}}}spPr")
    spPr.append(_solid_fill_elem(hex_color))
    return spPr

def _no_line_spPr(hex_color: str) -> etree._Element:
    """spPr with solid fill and no border line."""
    spPr = _spPr_with_fill(hex_color)
    ln   = etree.SubElement(spPr, f"{{{_A}}}ln")
    etree.SubElement(ln, f"{{{_A}}}noFill")
    return spPr

def _axis_txPr(hex_color: str, font_size_pt: int = 9) -> etree._Element:
    """Return <c:txPr> with white text for axis labels."""
    txPr  = etree.Element(f"{{{_C}}}txPr")
    bodyPr = etree.SubElement(txPr, f"{{{_A}}}bodyPr")
    lstStyle = etree.SubElement(txPr, f"{{{_A}}}lstStyle")
    p     = etree.SubElement(txPr, f"{{{_A}}}p")
    pPr   = etree.SubElement(p, f"{{{_A}}}pPr")
    defRPr = etree.SubElement(pPr, f"{{{_A}}}defRPr")
    defRPr.set("sz", str(font_size_pt * 100))
    defRPr.set("b",  "0")
    solidFill = _solid_fill_elem(hex_color)
    defRPr.append(solidFill)
    latin = etree.SubElement(defRPr, f"{{{_A}}}latin")
    latin.set("typeface", FONT_BODY)
    return txPr

def _set_element_fill(parent: etree._Element, hex_color: str):
    """Insert or replace <a:solidFill> inside parent's <c:spPr>."""
    spPr = parent.find(f"{{{_C}}}spPr")
    if spPr is None:
        spPr = etree.SubElement(parent, f"{{{_C}}}spPr")
    # Remove existing fills
    for old in spPr.findall(f"{{{_A}}}solidFill"):
        spPr.remove(old)
    for old in spPr.findall(f"{{{_A}}}noFill"):
        spPr.remove(old)
    spPr.insert(0, _solid_fill_elem(hex_color))

def _gridline_xml(hex_color: str, alpha: int = 80) -> etree._Element:
    """Create a subtle gridline spPr."""
    spPr = etree.Element(f"{{{_C}}}spPr")
    ln   = etree.SubElement(spPr, f"{{{_A}}}ln")
    fill = etree.SubElement(ln, f"{{{_A}}}solidFill")
    clr  = etree.SubElement(fill, f"{{{_A}}}srgbClr")
    clr.set("val", hex_color)
    alpha_elem = etree.SubElement(clr, f"{{{_A}}}alpha")
    alpha_elem.set("val", str(alpha * 1000))
    return spPr


def _replace_child(parent, new_elem):
    """Remove existing child with same tag as new_elem, then append new_elem."""
    tag = new_elem.tag
    old = parent.find(tag)
    if old is not None:
        parent.remove(old)
    parent.append(new_elem)


def style_chart_dark(chart, series_hex_colors: list[str]):
    """
    Apply dark navy theme to a pptx chart via direct XML manipulation.
    In python-pptx 1.x the chart XML part root is <c:chartSpace>;
    chart._element is <c:chart> (child of chartSpace).
    """
    # chart._element = <c:chart>
    # chartSpace     = <c:chartSpace>  (root of the chart XML part)
    chart_elem = chart._element
    chartSpace = chart_elem.getroottree().getroot()

    # 1. chartSpace background
    _replace_child(chartSpace, _no_line_spPr(BG_HEX))

    # 2. plotArea background
    plotArea = chart_elem.find(f"{{{_C}}}plotArea")
    if plotArea is not None:
        _replace_child(plotArea, _no_line_spPr(PANEL_HEX))

    # 3. Series fills
    for i, series in enumerate(chart.series):
        hex_c    = series_hex_colors[i % len(series_hex_colors)]
        ser_elem = series._element
        _replace_child(ser_elem, _no_line_spPr(hex_c))

    # 4. Category axis (x) — text colour
    try:
        if plotArea is not None:
            catAx = plotArea.find(f"{{{_C}}}catAx")
            if catAx is not None:
                _replace_child(catAx, _axis_txPr("8BA9C8"))
    except Exception:
        pass

    # 5. Value axis (y) — text colour + subtle gridlines
    try:
        if plotArea is not None:
            valAx = plotArea.find(f"{{{_C}}}valAx")
            if valAx is not None:
                _replace_child(valAx, _axis_txPr("8BA9C8"))
                gl = valAx.find(f"{{{_C}}}majorGridlines")
                if gl is not None:
                    _replace_child(gl, _gridline_xml("2C4A6E", alpha=60))
    except Exception:
        pass

    # 6. Legend text colour
    try:
        legend = chart_elem.find(f"{{{_C}}}legend")
        if legend is not None:
            _replace_child(legend, _axis_txPr("8BA9C8", font_size_pt=9))
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Presentation helpers
# ---------------------------------------------------------------------------

def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    """Add a completely blank slide (no placeholder layout)."""
    blank_layout = prs.slide_layouts[6]   # index 6 = blank
    return prs.slides.add_slide(blank_layout)


def set_slide_background(slide, hex_color: str = BG_HEX):
    """Fill slide background with a solid colour."""
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def add_textbox(slide, text: str, l, t, w, h,
                font_size=12, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font=FONT_BODY):
    txb  = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf   = txb.text_frame
    tf.word_wrap = False
    p    = tf.paragraphs[0]
    p.alignment = align
    run  = p.add_run()
    run.text = text
    run.font.name   = font
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_accent_band(slide):
    """Thin horizontal accent line below the title."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(BAND_L), Inches(BAND_T), Inches(BAND_W), Inches(BAND_H)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_footer(slide, left_text: str, right_text: str = ""):
    add_textbox(slide, left_text,
                FOOT_L, FOOT_T, FOOT_W / 2, FOOT_H,
                font_size=8, color=MUTED)
    if right_text:
        add_textbox(slide, right_text,
                    FOOT_L + FOOT_W / 2, FOOT_T, FOOT_W / 2, FOOT_H,
                    font_size=8, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Table builder  (dark-themed)
# ---------------------------------------------------------------------------

def add_dark_table(slide, df: pd.DataFrame,
                   money_cols: list[str] = None,
                   pct_cols:   list[str] = None,
                   col_widths: list[float] = None):
    """
    Add a dark-themed table to the fixed TABLE zone.
    money_cols → right-aligned with $ prefix and comma format
    pct_cols   → right-aligned with % suffix
    """
    money_cols = money_cols or []
    pct_cols   = pct_cols   or []

    nrows = len(df) + 1   # +1 header
    ncols = len(df.columns)

    tbl_shape = slide.shapes.add_table(
        nrows, ncols,
        Inches(TABLE_L), Inches(TABLE_T),
        Inches(TABLE_W), Inches(TABLE_H),
    )
    tbl = tbl_shape.table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    # Row heights
    tbl.rows[0].height = Inches(0.42)   # header
    for ri in range(1, nrows):
        tbl.rows[ri].height = Inches(0.56)

    def _set_cell(cell, text, bg_rgb, fg_rgb,
                  bold=False, font_size=10, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_rgb
        tf = cell.text_frame
        tf.word_wrap = False
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name  = FONT_BODY
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = fg_rgb

    # Header row
    for ci, col in enumerate(df.columns):
        label = col.replace("_", " ")
        _set_cell(
            tbl.cell(0, ci), label,
            bg_rgb=HDR, fg_rgb=WHITE,
            bold=True, font_size=10,
            align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT,
        )

    # Data rows
    for ri, (_, row_data) in enumerate(df.iterrows()):
        excel_ri = ri + 1
        alt_bg   = ALT if ri % 2 == 1 else PANEL

        for ci, (col, val) in enumerate(zip(df.columns, row_data)):
            if col in money_cols:
                display = f"${val:,.1f}"
                align   = PP_ALIGN.RIGHT
            elif col in pct_cols:
                display = f"{val:.1f}%"
                align   = PP_ALIGN.RIGHT
            elif col == "Quarter":
                display = str(val)
                align   = PP_ALIGN.LEFT
            else:
                display = str(val)
                align   = PP_ALIGN.CENTER

            _set_cell(
                tbl.cell(excel_ri, ci), display,
                bg_rgb=alt_bg,
                fg_rgb=WHITE if ci == 0 else MUTED,
                bold=(ci == 0),
                font_size=9,
                align=align,
            )

    return tbl_shape


# ---------------------------------------------------------------------------
# Chart builders
# ---------------------------------------------------------------------------

def add_clustered_bar(slide, df: pd.DataFrame, series_cols: list[str],
                      title: str = ""):
    """Clustered column chart — Revenue slide."""
    cd = ChartData()
    cd.categories = df["Quarter"].tolist()
    for col in series_cols:
        cd.add_series(col.replace("_", " "), df[col].tolist())

    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(CHART_L), Inches(CHART_T),
        Inches(CHART_W), Inches(CHART_H),
        cd,
    )
    chart = frame.chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    chart.has_legend = True
    chart.legend.position     = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    style_chart_dark(chart, PALETTE_HEX)
    return chart


def add_stacked_bar(slide, df: pd.DataFrame, series_cols: list[str],
                    title: str = ""):
    """Stacked column chart — CapEx slide."""
    cd = ChartData()
    cd.categories = df["Quarter"].tolist()
    for col in series_cols:
        cd.add_series(col.replace("_", " "), df[col].tolist())

    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(CHART_L), Inches(CHART_T),
        Inches(CHART_W), Inches(CHART_H),
        cd,
    )
    chart = frame.chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    chart.has_legend = True
    chart.legend.position     = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    style_chart_dark(chart, PALETTE_HEX)
    return chart


def add_line_chart(slide, df: pd.DataFrame, series_cols: list[str],
                   title: str = ""):
    """Line chart — Debt & Tax slide."""
    cd = ChartData()
    cd.categories = df["Quarter"].tolist()
    for col in series_cols:
        cd.add_series(col.replace("_", " "), df[col].tolist())

    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(CHART_L), Inches(CHART_T),
        Inches(CHART_W), Inches(CHART_H),
        cd,
    )
    chart = frame.chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    chart.has_legend = True
    chart.legend.position     = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Make lines thicker and set marker styles
    for i, series in enumerate(chart.series):
        series.smooth = True
        series.format.line.color.rgb = PALETTE[i % len(PALETTE)]
        series.format.line.width     = Pt(2.0)

    style_chart_dark(chart, PALETTE_HEX)
    return chart


# ---------------------------------------------------------------------------
# Slide constructors
# ---------------------------------------------------------------------------

def build_cover_slide(prs: Presentation, company: str, period: str, subtitle: str):
    slide = blank_slide(prs)
    set_slide_background(slide)

    # Large vertical accent bar on left
    bar = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0), Inches(0.30), Inches(7.50)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Company name
    add_textbox(
        slide, company.upper(),
        0.60, 1.60, 12.0, 0.80,
        font_size=32, bold=True, color=WHITE,
        align=PP_ALIGN.LEFT,
    )

    # Title
    add_textbox(
        slide, period,
        0.60, 2.55, 12.0, 1.10,
        font_size=40, bold=True, color=ACCENT,
        align=PP_ALIGN.LEFT,
    )

    # Subtitle
    add_textbox(
        slide, subtitle,
        0.60, 3.75, 10.0, 0.60,
        font_size=16, bold=False, color=MUTED,
        align=PP_ALIGN.LEFT,
    )

    # Divider line
    div = slide.shapes.add_shape(
        1,
        Inches(0.60), Inches(4.50), Inches(11.73), Inches(0.03)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = RULE
    div.line.fill.background()

    # Footer text
    add_textbox(
        slide, "CONFIDENTIAL — FOR INTERNAL USE ONLY",
        0.60, 6.90, 12.0, 0.40,
        font_size=9, color=MUTED, italic=True,
        align=PP_ALIGN.LEFT,
    )


def build_revenue_slide(prs: Presentation, df: pd.DataFrame, period: str):
    slide = blank_slide(prs)
    set_slide_background(slide)

    # Section label
    add_textbox(slide, "REVENUE", 0.35, 0.22, 2.5, 0.30,
                font_size=9, color=ACCENT, bold=True)

    # Title
    add_textbox(slide, "Revenue by Product Line",
                0.35, 0.48, 12.0, 0.55,
                font_size=22, bold=True, color=WHITE)

    add_accent_band(slide)

    col_widths = [1.15, 1.15, 1.15, 1.15, 1.30]
    add_dark_table(
        slide, df,
        money_cols=["Product_A", "Product_B", "Product_C", "Total_Revenue"],
        col_widths=col_widths,
    )

    add_clustered_bar(
        slide, df,
        series_cols=["Product_A", "Product_B", "Product_C"],
        title="Revenue by Product ($M)",
    )

    add_footer(slide, f"Revenue Trend  |  {period}", "USD $M")


def build_capex_slide(prs: Presentation, df: pd.DataFrame, period: str):
    slide = blank_slide(prs)
    set_slide_background(slide)

    add_textbox(slide, "CAPITAL EXPENDITURE", 0.35, 0.22, 4.0, 0.30,
                font_size=9, color=ACCENT, bold=True)

    add_textbox(slide, "CapEx Summary — Maintenance vs Growth",
                0.35, 0.48, 12.0, 0.55,
                font_size=22, bold=True, color=WHITE)

    add_accent_band(slide)

    col_widths = [1.15, 1.55, 1.40, 1.30]
    add_dark_table(
        slide, df,
        money_cols=["Maintenance_CapEx", "Growth_CapEx", "Total_CapEx"],
        col_widths=col_widths,
    )

    add_stacked_bar(
        slide, df,
        series_cols=["Maintenance_CapEx", "Growth_CapEx"],
        title="CapEx Split ($M)",
    )

    add_footer(slide, f"Capital Expenditure  |  {period}", "USD $M")


def build_debt_tax_slide(prs: Presentation, df: pd.DataFrame, period: str):
    slide = blank_slide(prs)
    set_slide_background(slide)

    add_textbox(slide, "DEBT & TAX", 0.35, 0.22, 3.0, 0.30,
                font_size=9, color=ACCENT, bold=True)

    add_textbox(slide, "Debt Reduction & Tax Overview",
                0.35, 0.48, 12.0, 0.55,
                font_size=22, bold=True, color=WHITE)

    add_accent_band(slide)

    col_widths = [1.15, 1.35, 1.25, 1.20, 0.95]
    add_dark_table(
        slide, df,
        money_cols=["Interest_Expense", "Total_Debt", "Tax_Expense"],
        pct_cols=["Effective_Tax_Rate"],
        col_widths=col_widths,
    )

    add_line_chart(
        slide, df,
        series_cols=["Interest_Expense", "Tax_Expense", "Effective_Tax_Rate"],
        title="Interest, Tax Expense & Effective Rate",
    )

    add_footer(slide, f"Debt & Tax  |  {period}", "USD $M  |  Rate %")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def load_excel(path: Path) -> dict[str, pd.DataFrame]:
    xl = pd.ExcelFile(path)
    # Row 0 is the sheet title; actual column headers are in row 1.
    # Headers were written with spaces; convert back to underscores so the
    # rest of the script can reference columns by their original names.
    sheets = {}
    for sheet in xl.sheet_names:
        df = xl.parse(sheet, header=1)
        df.columns = [str(c).replace(" ", "_") for c in df.columns]
        sheets[sheet] = df
    return sheets


def main():
    parser = argparse.ArgumentParser(description="Excel → PowerPoint report generator")
    parser.add_argument("--input",   default="data/financial_report.xlsx")
    parser.add_argument("--output",  default="output/quarterly_report.pptx")
    parser.add_argument("--company", default="Acme Corporation")
    parser.add_argument("--period",  default="Q4 2024 Financial Review")
    args = parser.parse_args()

    input_path  = Path(args.input)
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    print(f"Reading:  {input_path}")
    sheets = load_excel(input_path)

    df_rev  = sheets["Revenue_Trend"]
    df_cap  = sheets["CapEx"]
    df_debt = sheets["Debt_and_Tax"]

    print("Building presentation...")
    prs = new_presentation()

    build_cover_slide(prs, args.company, args.period,
                      "Quarterly Financial Review — Management Pack")
    print("  [1/4] Cover slide")

    build_revenue_slide(prs, df_rev, args.period)
    print("  [2/4] Revenue slide")

    build_capex_slide(prs, df_cap, args.period)
    print("  [3/4] CapEx slide")

    build_debt_tax_slide(prs, df_debt, args.period)
    print("  [4/4] Debt & Tax slide")

    prs.save(output_path)
    print(f"\nSaved:  {output_path}")
    print(f"Slides: {len(prs.slides)}  |  Size: {output_path.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
